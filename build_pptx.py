#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

RED   = RGBColor(0xEB, 0x00, 0x00)
DARK  = RGBColor(0x1C, 0x1C, 0x1C)
GREY  = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xFB, 0xFA, 0xFA)
FONT  = "Liberation Sans"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf

def setpara(p, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
            space_after=6, font=FONT):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return p

def bullets(tf, items, size=16, color=DARK, gap=8):
    for it in items:
        # reuse the initial paragraph only if it is still empty, else append
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "•  " + it
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        r = p.runs[0]
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def heading(s, title, num=None):
    # top red bar
    rect(s, 0, 0, SW, Inches(1.15), RED)
    tf = textbox(s, Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.8),
                 anchor=MSO_ANCHOR.MIDDLE)
    label = (f"{num}.  " if num else "") + title
    setpara(tf.paragraphs[0], label, 26, WHITE, bold=True)

# ---------------------------------------------------------------- SLIDE 1 (titre)
s = slide()
rect(s, 0, 0, SW, SH, RED)
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
setpara(tf.paragraphs[0], "CHARTE ÉTHIQUE & DURABILITÉ NUMÉRIQUE", 15, WHITE,
        bold=True, space_after=14)
p = tf.add_paragraph(); setpara(p, "La numérisation des titres de transport aux CFF",
        40, WHITE, bold=True, space_after=10)
p = tf.add_paragraph(); setpara(p, "Fondements, engagements et réflexion critique",
        20, WHITE, italic=True, space_after=0)
tf2 = textbox(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(1.0))
setpara(tf2.paragraphs[0], "Cours Éthique et durabilité — HEIG-VD", 15, WHITE)
p = tf2.add_paragraph(); setpara(p, "Sofian Ethenoz  ·  Arthur Menétrey  ·  Mirco Profico", 15, WHITE)
notes(s,
"[VOIX 1 — Sofian]\n"
"Bonjour. Nous vous présentons notre charte éthique et de durabilité numérique, "
"construite autour d'un cas concret : la numérisation des titres de transport aux "
"CFF, l'entreprise de transport public de toute la Suisse. Notre fil conducteur est "
"simple : un outil pensé pour faciliter la mobilité du plus grand nombre ne doit "
"jamais finir par exclure celles et ceux qui dépendent le plus du service public.")

# ---------------------------------------------------------------- SLIDE 2 (entreprise)
s = slide()
heading(s, "Les CFF, leurs technologies et leurs bénéficiaires")
tfL = textbox(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.6))
setpara(tfL.paragraphs[0], "Des titres de transport numérisés", 18, RED, bold=True, space_after=10)
bullets(tfL, [
 "Application Mobile CFF — billets sur smartphone",
 "SwissPass — carte à puce (AG, demi-tarif…)",
 "EasyRide — achat automatique « check-in / check-out »",
 "Automates et vente en ligne",
], size=15)
tfR = textbox(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.6))
setpara(tfR.paragraphs[0], "Un service public : servir TOUT le monde", 18, RED, bold=True, space_after=10)
p = tfR.add_paragraph(); setpara(p, "Des personnes vulnérables face au numérique :", 15, DARK, space_after=8)
bullets(tfR, [
 "Personnes âgées (smartphone, paiement)",
 "Personnes en situation de handicap (accessibilité)",
 "Personnes à faibles revenus (pas de smartphone/data)",
 "Jeunes, mineurs et touristes",
], size=15)
notes(s,
"[VOIX 1 — Sofian]\n"
"Les CFF ont numérisé leurs titres : application Mobile CFF, carte SwissPass, système "
"EasyRide par géolocalisation, vente en ligne qui remplace les guichets. Mais un "
"service public doit servir tout le monde. Or, parmi les usagers, se trouvent des "
"personnes vulnérables — personnes âgées, en situation de handicap, à faibles revenus, "
"jeunes ou touristes — pour qui le tout-numérique peut devenir une barrière. C'est de "
"ce constat que naît notre charte.")

# ---------------------------------------------------------------- SLIDE 3 (axes 1-2)
s = slide()
heading(s, "Notre charte (1/3) — Autonomie & vulnérabilité")
tfL = textbox(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(5.7))
setpara(tfL.paragraphs[0], "Axe 1 — Finalité & autonomie", 18, RED, bold=True, space_after=8)
bullets(tfL, [
 "La mobilité reste un droit, pas un privilège des « connectés »",
 "Toujours une alternative non numérique (guichet, papier)",
 "Le numérique ne coûte jamais moins cher : pas de contrainte",
], size=15)
tfR = textbox(s, Inches(6.9), Inches(1.45), Inches(5.8), Inches(5.7))
setpara(tfR.paragraphs[0], "Axe 2 — Vulnérabilité & IA", 18, RED, bold=True, space_after=8)
bullets(tfR, [
 "Accessibilité & design universel (WCAG), conçus AVEC les concernés",
 "Aucune stigmatisation de qui n'utilise pas le numérique",
 "IA : transparence, « human in the loop », contre les biais",
 "L'IA assiste le jugement humain, ne le remplace pas",
], size=15)
notes(s,
"[VOIX 2 — Arthur]\n"
"Notre charte tient en treize engagements, regroupés en cinq axes. Premier axe, "
"l'autonomie : la mobilité reste un droit, pas un privilège des personnes connectées ; "
"pour chaque service numérique, une alternative — guichet ou billet papier — jamais "
"plus chère. Deuxième axe, la vulnérabilité et l'intelligence artificielle : des "
"interfaces accessibles, conçues avec les personnes concernées ; une IA transparente, "
"sous contrôle humain et auditée contre les biais — elle assiste le jugement, sans le "
"remplacer.")

# ---------------------------------------------------------------- SLIDE 4 (axes 3-4)
s = slide()
heading(s, "Notre charte (2/3) — Hybridation & valeurs")
tfL = textbox(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(5.7))
setpara(tfL.paragraphs[0], "Axe 3 — Hybridation au service de l'humain", 18, RED, bold=True, space_after=8)
bullets(tfL, [
 "La technologie augmente l'usager et libère le personnel",
 "La présence humaine en gare est un choix assumé",
 "Limite : pas de biométrie imposée, voyager anonyme reste possible",
], size=15)
tfR = textbox(s, Inches(6.9), Inches(1.45), Inches(5.8), Inches(5.7))
setpara(tfR.paragraphs[0], "Axe 4 — Valeurs & préjudices", 18, RED, bold=True, space_after=8)
bullets(tfR, [
 "Solidarité, égalité d'accès, confiance, vie privée",
 "Données minimisées, hébergées en Suisse, jamais revendues",
 "Refus d'un historique permanent des déplacements",
 "Suivi de l'inclusion : ne laisser personne de côté",
], size=15)
notes(s,
"[VOIX 2 — Arthur]\n"
"Troisième axe, l'hybridation entre humain et technologie : nous la voulons au service "
"de l'humain, avec une présence humaine assumée en gare. Notre limite : pas de "
"biométrie imposée, et il reste possible de voyager anonymement. Quatrième axe, les "
"valeurs : solidarité, égalité d'accès, confiance et vie privée. Les données sont "
"minimisées, hébergées en Suisse, jamais revendues ; nous refusons tout historique "
"permanent des déplacements.")

# ---------------------------------------------------------------- SLIDE 5 (axe 5 durabilité)
s = slide()
heading(s, "Notre charte (3/3) — Durabilité & sobriété")
tfL = textbox(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(5.7))
setpara(tfL.paragraphs[0], "Axe 5 — Sobriété numérique", 18, RED, bold=True, space_after=8)
bullets(tfL, [
 "Numériser ce qui a une valeur réelle, pas multiplier les usages",
 "Apps légères pour smartphones anciens — anti-obsolescence",
 "Supports durables (SwissPass), data centers renouvelables",
], size=15)
tfR = textbox(s, Inches(6.9), Inches(1.45), Inches(5.8), Inches(5.7))
setpara(tfR.paragraphs[0], "Bénéfice net, pas seulement « moins de papier »", 18, RED, bold=True, space_after=8)
bullets(tfR, [
 "Lecture critique : le numérique n'est pas immatériel",
 "Risque d'effet rebond (fabrication des smartphones)",
 "Mesurer et publier l'empreinte réelle",
 "Évaluer la charte chaque année avec les associations",
], size=15)
notes(s,
"[VOIX 3 — Mirco]\n"
"Cinquième axe : la durabilité. Nous inscrivons le numérique dans une logique de "
"sobriété — numériser ce qui a une vraie valeur, sans multiplier les usages. Des "
"applications légères, compatibles avec les anciens smartphones pour éviter "
"l'obsolescence, des data centers à l'électricité renouvelable. Et une lecture "
"critique : le numérique n'est pas immatériel. Supprimer le papier ne suffit pas si "
"l'on pousse à changer de téléphone — c'est l'effet rebond. Nous visons donc un "
"bénéfice net, mesuré et publié.")

# ---------------------------------------------------------------- SLIDE 6 (conclusion)
s = slide()
rect(s, 0, 0, SW, SH, RED)
tf = textbox(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(4.5))
setpara(tf.paragraphs[0], "Notre conviction", 18, WHITE, bold=True, space_after=18)
p = tf.add_paragraph(); setpara(p,
 "Nos choix de conception ne sont jamais neutres :", 26, WHITE, bold=True, space_after=12)
p = tf.add_paragraph(); setpara(p,
 "une décision technique est souvent une décision éthique déguisée.", 26, WHITE, bold=True, space_after=24)
p = tf.add_paragraph(); setpara(p,
 "Concevoir une mobilité numérique qui n'oublie personne —", 20, WHITE, italic=True, space_after=4)
p = tf.add_paragraph(); setpara(p,
 "sobre, accessible et respectueuse de la vie privée.", 20, WHITE, italic=True, space_after=0)
tf2 = textbox(s, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.8))
setpara(tf2.paragraphs[0], "Merci de votre attention.", 18, WHITE, bold=True)
notes(s,
"[VOIX 3 — Mirco]\n"
"En conclusion, nos choix de conception ne sont jamais neutres : qu'une application "
"tourne sur un vieux téléphone, qu'un guichet reste ouvert ou qu'un algorithme n'ait "
"pas le dernier mot, ce sont des décisions éthiques déguisées en décisions techniques. "
"Nous voulons intégrer l'éthique et la durabilité dès le départ, et garder les plus "
"vulnérables comme mesure de la qualité d'une technologie. Concevoir une mobilité "
"numérique qui n'oublie personne : voilà notre engagement. Merci de votre attention.")

# --- clean document metadata (no tool traces) ---
cp = prs.core_properties
cp.title = "Charte éthique et durabilité numérique — Numérisation des titres de transport aux CFF"
cp.subject = "Éthique et durabilité — HEIG-VD"
cp.author = "Groupe Éthique et durabilité"
cp.last_modified_by = ""
cp.comments = ""
cp.category = ""
cp.keywords = ""

prs.save("presentation_charte_CFF.pptx")
print("OK presentation_charte_CFF.pptx  —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
